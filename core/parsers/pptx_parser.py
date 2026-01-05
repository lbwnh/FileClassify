"""
PPTX 解析器实现。

本模块提供 PowerPoint 文件解析功能。
"""

from .base import BaseParser
from typing import Dict


class PPTXParser(BaseParser):
    """
    PowerPoint 文件解析器（PPTX, PPT）。

    使用 python-pptx 进行文本提取。
    """
    
    def extract_text(self) -> str:
        """
        从 PowerPoint 文件中提取文本内容。

        返回:
            str: 提取的文本内容
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"[幻灯片 {slide_num}]")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                    
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = '\t'.join([cell.text for cell in row.cells])
                            if row_text.strip():
                                text_content.append(row_text)
            
            return '\n'.join(text_content)
        
        except ImportError:
            return f"[PPTX Parser] 无法提取文本：python-pptx 未安装。\n文件：{self.file_path.name}"
        
        except Exception as e:
            return f"[PPTX Parser] 提取文本出错：{str(e)}\n文件：{self.file_path.name}"
    
    def extract_metadata(self) -> Dict[str, any]:
        """
        从 PowerPoint 文件中提取元数据。

        返回:
            dict: PowerPoint 元数据
        """
        metadata = {
            'title': None,
            'author': None,
            'subject': None,
            'keywords': None,
            'created_date': None,
            'modified_date': None,
            'last_modified_by': None,
            'slide_count': 0
        }
        
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            
            core_props = prs.core_properties
            
            metadata['title'] = core_props.title
            metadata['author'] = core_props.author
            metadata['subject'] = core_props.subject
            metadata['keywords'] = core_props.keywords
            metadata['created'] = core_props.created
            metadata['modified'] = core_props.modified
            metadata['last_modified_by'] = core_props.last_modified_by
            
            metadata['slide_count'] = len(prs.slides)
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata
